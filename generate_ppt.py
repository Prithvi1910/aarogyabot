from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AarogyaBot (आरोग्यबॉट)"
subtitle.text = "Multilingual Healthcare Assistant for Rural India\n\nVision: Bridging the healthcare gap by providing instant, accurate, and localized medical assessments to communities with limited access to immediate medical care."

# Slide 2: The Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "The Problem"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "1. Doctor Shortage: Rural India faces a severe shortage of specialized medical professionals."
tf.add_paragraph().text = "2. Language Barriers: Most medical information is in English."
tf.add_paragraph().text = "3. Misinformation: Reliance on inaccurate home remedies for severe diseases."
tf.add_paragraph().text = "4. Delayed Treatment: Patients wait until symptoms become critical."

# Slide 3: The Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "The Solution: AarogyaBot"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "AarogyaBot acts as a Top-Level Expert Doctor accessible via smartphone."
tf.add_paragraph().text = "- Multilingual: Hindi, Tamil, Telugu, Gujarati, Marathi, Hinglish."
tf.add_paragraph().text = "- Voice-Enabled: Uses native Speech-to-Text."
tf.add_paragraph().text = "- Expert Guidance: Diagnoses symptom clusters & suggests medicine."
tf.add_paragraph().text = "- Facility Locator: Instantly locates nearest PHCs using a PIN code."

# Slide 4: Technical Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Technical Architecture"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "Modern decoupled Client-Server architecture:"
tf.add_paragraph().text = "- Frontend: React, Vite, TailwindCSS, Web Speech API"
tf.add_paragraph().text = "- Backend: FastAPI, Python"
tf.add_paragraph().text = "- AI Core: LangChain, Groq API (Llama-3.1), FAISS Vector DB"
tf.add_paragraph().text = "- Translation: Google Translate (deep-translator)"

# Slide 5: How the AI Works (RAG + Memory)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "How the AI Works (RAG + Memory)"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "1. Knowledge Base: Strict database of rural diseases and treatments."
tf.add_paragraph().text = "2. Vector Retrieval: FAISS matches user symptoms against knowledge base."
tf.add_paragraph().text = "3. Conversational Memory: Tracks session history for natural follow-ups."
tf.add_paragraph().text = "4. Safety Guardrails: Emergency protocols and strict reliance on provided context."

# Slide 6: Future Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Roadmap"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "1. Text-to-Speech (TTS): Read medical advice out loud."
tf.add_paragraph().text = "2. WhatsApp Integration: Deploy via Twilio API."
tf.add_paragraph().text = "3. Image Recognition: Diagnose visual symptoms (rashes, wounds)."
tf.add_paragraph().text = "4. Persistent Database: Move memory from RAM to Redis/Supabase."

prs.save("C:\\Users\\vikra\\aarogyabot\\presentation_deck.pptx")
print("Presentation successfully created!")
